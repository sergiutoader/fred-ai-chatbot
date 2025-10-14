from pptx import Presentation
from pptx.util import Pt

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Demo Slide"
body = slide.placeholders[1]
body.name = "FRED_CONTENT"
body.text = "Replace me with LLM output"
for p in body.text_frame.paragraphs: p.font.size = Pt(18)
prs.save("simple_template.pptx")
print("✅ Created simple_template.pptx in current directory")
